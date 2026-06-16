"""Document upload + library service for RAG.

Supported input formats:
  - text/plain, text/markdown          (decode bytes)
  - application/pdf                    (pypdf — pure-python, no ML)
  - .docx (Word)                       (python-docx)
  - .xlsx (Excel)                      (openpyxl — sheet+cell flatten)
  - .pptx (PowerPoint)                 (python-pptx — slide text flatten)

We extract text at upload time, never re-parse the original. Storage is
just the extracted text + the original filename for display. We DON'T
store the raw bytes — that would 5–20× the storage cost for zero value.

Why these parsers and not docling? Docling does layout-aware extraction
with IBM ML models (~1GB venv, 5-30s/page on CPU). For prose-heavy chat
uploads (the 95% case) deterministic parsers are faster and cheaper.
If a user shows up with complex-table PDFs, we can add docling as an
opt-in `advanced=true` flag without rewriting this dispatcher.

Limits:
  - max file size 5 MB before extraction
  - max extracted text 200_000 chars after extraction (~50 pages)
  - per-user document quota 50

RAG v2 chunks + embeds for retrieval; if indexing fails the doc still
works in v1's "prepend full text" mode.
"""
from __future__ import annotations

import io
import logging
import uuid
from typing import Optional

from sqlalchemy.orm import Session

from app.models.document import KaiDocument

logger = logging.getLogger(__name__)

MAX_FILE_BYTES = 5 * 1024 * 1024       # 5 MB raw
MAX_TEXT_CHARS = 200_000               # ~50 pages
MAX_DOCS_PER_USER = 50


class DocumentError(ValueError):
    """User-facing error: surfaces directly to API 4xx response."""


def _extract_pdf(data: bytes) -> str:
    # pypdf is the maintained successor to PyPDF2 — same API style,
    # better tolerance for malformed PDFs, faster on large files.
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover — should be in venv
        raise DocumentError("PDF support not installed on server") from e
    try:
        reader = PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            t = page.extract_text() or ""
            if t.strip():
                parts.append(t.strip())
        return "\n\n".join(parts)
    except Exception as e:
        raise DocumentError(f"could not parse PDF: {e}") from e


def _extract_docx(data: bytes) -> str:
    try:
        import docx  # python-docx
    except ImportError as e:  # pragma: no cover
        raise DocumentError("Word support not installed on server") from e
    try:
        d = docx.Document(io.BytesIO(data))
    except Exception as e:
        raise DocumentError(f"could not parse Word doc: {e}") from e
    parts: list[str] = []
    # Paragraphs in body order
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    # Tables flattened row-by-row as tab-separated cells — keeps reading
    # order coherent in the extracted text.
    for table in d.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            line = "\t".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    try:
        import openpyxl
    except ImportError as e:  # pragma: no cover
        raise DocumentError("Excel support not installed on server") from e
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception as e:
        raise DocumentError(f"could not parse Excel workbook: {e}") from e
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            line = "\t".join(cells).rstrip("\t")
            if line.strip():
                parts.append(line)
    wb.close()
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise DocumentError("PowerPoint support not installed on server") from e
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise DocumentError(f"could not parse PowerPoint deck: {e}") from e
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if t:
                parts.append(t)
    return "\n".join(parts)


def _extract_text(data: bytes) -> str:
    # Best-effort decode — try UTF-8 first, then latin-1 as fallback.
    for enc in ("utf-8", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    raise DocumentError("could not decode text file (not UTF-8 or latin-1)")


def upload(
    db: Session,
    user_id: uuid.UUID,
    filename: str,
    mime_type: str,
    data: bytes,
) -> KaiDocument:
    """Extract text from data + persist. Returns the row.

    Raises DocumentError on any user-visible failure (too big, bad PDF,
    bad encoding, quota exceeded). The caller maps to HTTPException(400/413).
    """
    if not filename or not filename.strip():
        raise DocumentError("filename required")
    if len(data) > MAX_FILE_BYTES:
        raise DocumentError(
            f"file too large ({len(data):,} bytes); max {MAX_FILE_BYTES:,}"
        )
    if not data:
        raise DocumentError("empty file")

    mime = (mime_type or "").lower().strip()
    fname_lower = filename.lower()
    # Dispatch on extension when available (more reliable than MIME — browsers
    # often send application/octet-stream for .docx/.xlsx/.pptx) with a MIME
    # fallback for clients that do send the right Content-Type.
    if mime == "application/pdf" or fname_lower.endswith(".pdf"):
        text = _extract_pdf(data)
        mime = "application/pdf"
    elif fname_lower.endswith(".docx") or mime == (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        text = _extract_docx(data)
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif fname_lower.endswith(".xlsx") or mime == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        text = _extract_xlsx(data)
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif fname_lower.endswith(".pptx") or mime == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        text = _extract_pptx(data)
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif mime.startswith("text/") or fname_lower.endswith((".txt", ".md", ".markdown")):
        text = _extract_text(data)
        if not mime:
            mime = "text/plain"
    else:
        raise DocumentError(
            f"unsupported file type {mime!r}; supported: PDF, Word (.docx), "
            f"Excel (.xlsx), PowerPoint (.pptx), text, markdown"
        )

    text = (text or "").strip()
    if not text:
        raise DocumentError("no extractable text in file")
    if len(text) > MAX_TEXT_CHARS:
        # Truncate rather than reject — most over-limit docs are slightly over
        text = text[:MAX_TEXT_CHARS]
        logger.warning("document truncated to %d chars (was over limit)", MAX_TEXT_CHARS)

    # Quota check
    current = db.query(KaiDocument).filter(KaiDocument.user_id == user_id).count()
    if current >= MAX_DOCS_PER_USER:
        raise DocumentError(
            f"document quota reached ({MAX_DOCS_PER_USER}); delete one first"
        )

    row = KaiDocument(
        user_id=user_id,
        filename=filename.strip()[:200],  # cap display name
        mime_type=mime,
        text_len=len(text),
        full_text=text,
    )
    db.add(row)
    db.commit()
    db.refresh(row)

    # RAG v2: chunk + embed for retrieval. Failure here is non-fatal — the
    # doc still works in v1's "prepend full text" mode if no chunks index.
    try:
        from app.services import rag
        n_chunks = rag.index_document(db, row.id, user_id, text)
        logger.info("indexed %d chunks for doc %s", n_chunks, row.id)
    except Exception as e:
        logger.warning("chunk indexing failed for doc %s: %s", row.id, e)

    return row


def list_user(db: Session, user_id: uuid.UUID) -> list[KaiDocument]:
    return (
        db.query(KaiDocument)
        .filter(KaiDocument.user_id == user_id)
        .order_by(KaiDocument.created_at.desc())
        .all()
    )


def get_for_user(
    db: Session, user_id: uuid.UUID, doc_id: uuid.UUID
) -> Optional[KaiDocument]:
    return (
        db.query(KaiDocument)
        .filter(KaiDocument.id == doc_id, KaiDocument.user_id == user_id)
        .first()
    )


def delete(db: Session, user_id: uuid.UUID, doc_id: uuid.UUID) -> bool:
    row = get_for_user(db, user_id, doc_id)
    if row is None:
        return False
    db.delete(row)
    db.commit()
    return True
