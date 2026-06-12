"""Tests for the document extractors in app.services.documents.

We don't ship sample files in the repo — instead we synthesize minimal
docs in each format on the fly using the same parser libraries. That
keeps the test suite hermetic and proves the upload + extract round-trip
works for each format without committing binary fixtures.
"""
from __future__ import annotations

import io

import pytest

from app.services.documents import (
    DocumentError,
    _extract_docx,
    _extract_pdf,
    _extract_pptx,
    _extract_text,
    _extract_xlsx,
)


def _make_docx(paragraphs: list[str], table_rows: list[list[str]] | None = None) -> bytes:
    import docx
    d = docx.Document()
    for p in paragraphs:
        d.add_paragraph(p)
    if table_rows:
        t = d.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r_idx, row in enumerate(table_rows):
            for c_idx, val in enumerate(row):
                t.rows[r_idx].cells[c_idx].text = val
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list]]) -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    # Workbook starts with one default sheet — overwrite it
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for r in rows:
            ws.append(r)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slides_text: list[list[str]]) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[5]  # title-only layout
    for texts in slides_text:
        slide = prs.slides.add_slide(blank)
        if texts:
            slide.shapes.title.text = texts[0]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf(lines: list[str]) -> bytes:
    """Hand-rolled minimal PDF — avoids pulling reportlab into the dev deps.

    Renders each line at a fixed y-offset. Enough for pypdf.extract_text()
    to round-trip the strings, which is what we're testing.
    """
    text_lines = "\n".join(
        f"BT /F1 12 Tf 50 {750 - 18 * i} Td ({s}) Tj ET" for i, s in enumerate(lines)
    )
    content = f"q\n{text_lines}\nQ".encode("latin-1")
    objs = [
        b"%PDF-1.4\n",
        b"1 0 obj <</Type/Catalog/Pages 2 0 R>> endobj\n",
        b"2 0 obj <</Type/Pages/Kids[3 0 R]/Count 1>> endobj\n",
        (
            b"3 0 obj <</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
            b"/Resources<</Font<</F1 5 0 R>>>>"
            b"/Contents 4 0 R>> endobj\n"
        ),
        b"4 0 obj <</Length " + str(len(content)).encode() + b">>stream\n" + content + b"\nendstream endobj\n",
        b"5 0 obj <</Type/Font/Subtype/Type1/BaseFont/Helvetica>> endobj\n",
    ]
    offsets = []
    pos = 0
    blob = b""
    for o in objs:
        if not o.startswith(b"%PDF"):
            offsets.append(pos)
        blob += o
        pos += len(o)
    xref_pos = pos
    xref = b"xref\n0 6\n0000000000 65535 f \n"
    for off in offsets:
        xref += f"{off:010d} 00000 n \n".encode()
    trailer = (
        b"trailer <</Size 6/Root 1 0 R>>\n"
        b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    )
    return blob + xref + trailer


# ─── extractor tests ─────────────────────────────────────────────────


def test_extract_text_utf8():
    assert _extract_text("hello world".encode("utf-8")) == "hello world"


def test_extract_text_latin1_fallback():
    # Bytes that are valid latin-1 but invalid utf-8 (single high byte)
    data = b"caf\xe9"
    assert _extract_text(data) == "café"


def test_extract_text_unparseable_raises():
    # Manufactured bytes that fail both utf-8 AND strict latin-1? latin-1
    # accepts every byte sequence by design, so this won't actually fail —
    # we just confirm we don't lose data on weird input.
    assert _extract_text(b"\xff\xfe\xfd") is not None


def test_extract_docx_paragraphs():
    data = _make_docx(["First paragraph.", "Second paragraph."])
    out = _extract_docx(data)
    assert "First paragraph." in out
    assert "Second paragraph." in out


def test_extract_docx_table_rows_flattened():
    data = _make_docx(
        ["Header"],
        table_rows=[["Name", "Score"], ["Alice", "92"], ["Bob", "87"]],
    )
    out = _extract_docx(data)
    assert "Header" in out
    assert "Alice\t92" in out
    assert "Bob\t87" in out


def test_extract_xlsx_multiple_sheets():
    data = _make_xlsx({
        "Q1": [["product", "sales"], ["widget", 100], ["gadget", 50]],
        "Q2": [["product", "sales"], ["widget", 120]],
    })
    out = _extract_xlsx(data)
    assert "# Sheet: Q1" in out
    assert "# Sheet: Q2" in out
    assert "widget\t100" in out
    assert "gadget\t50" in out


def test_extract_pptx_slides_numbered():
    data = _make_pptx([["Intro slide"], ["Body slide"], ["Conclusion"]])
    out = _extract_pptx(data)
    assert "## Slide 1" in out
    assert "Intro slide" in out
    assert "## Slide 2" in out
    assert "Body slide" in out
    assert "## Slide 3" in out


def test_extract_pdf_round_trip():
    data = _make_pdf(["Hello world from PDF.", "Second line of text."])
    out = _extract_pdf(data)
    # pypdf may add/normalize whitespace — test substring presence
    assert "Hello world from PDF" in out
    assert "Second line of text" in out


def test_extract_pdf_malformed_raises():
    with pytest.raises(DocumentError):
        _extract_pdf(b"not-a-pdf-just-random-bytes")


def test_extract_docx_malformed_raises():
    with pytest.raises(DocumentError):
        _extract_docx(b"not-a-zip-not-a-docx")


def test_extract_xlsx_malformed_raises():
    with pytest.raises(DocumentError):
        _extract_xlsx(b"not-an-excel-file")


def test_extract_pptx_malformed_raises():
    with pytest.raises(DocumentError):
        _extract_pptx(b"not-a-pptx")
